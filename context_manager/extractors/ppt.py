"""PowerPoint extractor: per-slide title, body text, and speaker notes."""
from __future__ import annotations

from ..models import ExtractedDocument, ExtractedSection
from .base import ExtractionError


def extract(path, max_chars_per_slide: int = 2000) -> ExtractedDocument:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ExtractionError("python-pptx is required to read PowerPoint files") from e

    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise ExtractionError(f"failed to open presentation: {e}") from e

    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        title = ""
        if title_shape is not None and title_shape.has_text_frame:
            title = title_shape.text_frame.text.strip()

        texts = []
        for shape in slide.shapes:
            if shape is title_shape or not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)

        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        content = f"title: {title}\n" + "\n".join(texts)
        if notes:
            content += f"\nnotes: {notes}"
        if len(content) > max_chars_per_slide:
            content = content[:max_chars_per_slide] + "\n...(truncated)"

        location = f"Slide {i}: {title}" if title else f"Slide {i}"
        sections.append(ExtractedSection(location=location, text=content))

    return ExtractedDocument(kind="pptx", sections=sections)
