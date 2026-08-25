"""Reader: on-demand, read-only access to the original files.

Used by an Agent after search_context() has narrowed down a candidate
resource + location. Reads more completely than the Extractor's indexing-time
peek, but is still bounded (max_chars / max_rows) - narrow with `location`,
or use the Excel-specific range/formula helpers, instead of pulling an
entire huge file in one call. Never writes to the original file.
"""
from __future__ import annotations

from pathlib import Path

from . import extractors
from .extractors import excel
from .extractors.base import ExtractionError

EXCEL_EXTENSIONS = {"xlsx", "xlsm"}


def _ext(path) -> str:
    return Path(path).suffix.lower().lstrip(".")


def list_resource_structure(path: str) -> dict:
    """List the locations (sheets/slides/sections/pages) a resource
    contains, without reading their full content."""
    doc = extractors.extract_file(path)
    return {"kind": doc.kind, "locations": [s.location for s in doc.sections]}


def read_resource(path: str, location: str | None = None, max_chars: int = 12000) -> str:
    """Read a resource's content, optionally scoped to one location."""
    ext = _ext(path)
    if ext in EXCEL_EXTENSIONS:
        return _read_excel(path, location, max_chars)
    if ext == "pptx":
        return _read_pptx(path, location, max_chars)
    if ext == "pdf":
        return _read_pdf(path, location, max_chars)
    if ext == "docx":
        return _read_docx(path, location, max_chars)
    if ext in ("csv", "txt", "md"):
        return _read_text(path, max_chars)
    raise ExtractionError(f"unsupported extension: {ext}")


def _truncate(text: str, max_chars: int) -> str:
    if len(text) > max_chars:
        return text[:max_chars] + "\n...(truncated - narrow the location or use a range read)"
    return text


def _read_text(path, max_chars: int) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return _truncate(f.read(max_chars + 1), max_chars)


def _read_excel(path, location, max_chars: int) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        names = [location] if location and location in wb.sheetnames else wb.sheetnames
        blocks = []
        for name in names:
            ws = wb[name]
            lines = [f"[{name}] used_range={excel.used_range(ws)}"]
            try:
                merged = list(ws.merged_cells.ranges)
            except Exception:
                merged = []
            if merged:
                lines.append("merged_cells: " + ", ".join(str(m) for m in merged[:50]))
            for row in ws.iter_rows(max_row=300, max_col=40, values_only=True):
                if any(c is not None for c in row):
                    lines.append(" | ".join("" if c is None else str(c) for c in row))
            blocks.append("\n".join(lines))
        return _truncate("\n\n".join(blocks), max_chars)
    finally:
        wb.close()


def _read_pptx(path, location, max_chars: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        label = f"Slide {i}"
        if location and location != label and not location.startswith(f"{label}:"):
            continue
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        block = f"[{label}]\n" + "\n".join(texts)
        if notes:
            block += f"\nnotes: {notes}"
        blocks.append(block)
    return _truncate("\n\n".join(blocks), max_chars)


def _read_pdf(path, location, max_chars: int) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    blocks = []
    for i, page in enumerate(reader.pages, start=1):
        label = f"Page {i}"
        if location and location != label:
            continue
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        blocks.append(f"[{label}]\n{text}")
    return _truncate("\n\n".join(blocks), max_chars)


def _read_docx(path, location, max_chars: int) -> str:
    from docx import Document

    doc = Document(str(path))
    blocks: list[str] = []
    current_heading = "Introduction"
    current_lines: list[str] = []

    def flush() -> None:
        if current_lines and (location is None or location == current_heading):
            blocks.append(f"[{current_heading}]\n" + "\n".join(current_lines))

    for para in doc.paragraphs:
        style = (para.style.name or "") if para.style else ""
        if style.startswith("Heading") or style == "Title":
            flush()
            current_heading = para.text.strip() or current_heading
            current_lines = []
        elif para.text.strip():
            current_lines.append(para.text)
    flush()

    if not blocks:
        blocks.append("\n".join(p.text for p in doc.paragraphs if p.text.strip()))
    return _truncate("\n\n".join(blocks), max_chars)


# ---- Excel-specific deep-read helpers -----------------------------------------


def list_sheets(path: str) -> list[str]:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True)
    try:
        return wb.sheetnames
    finally:
        wb.close()


def read_range(path: str, sheet: str, cell_range: str, max_rows: int = 500, max_cols: int = 100) -> list[list]:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        ws = wb[sheet]
        selection = ws[cell_range]
        if hasattr(selection, "value"):  # single cell
            return [[selection.value]]
        rows = []
        for i, row in enumerate(selection):
            if i >= max_rows:
                break
            row_iter = row if isinstance(row, tuple) else (row,)
            rows.append([c.value for c in row_iter[:max_cols]])
        return rows
    finally:
        wb.close()


def read_formula(path: str, sheet: str, cell: str) -> str | None:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=False)
    try:
        ws = wb[sheet]
        value = ws[cell].value
        if isinstance(value, str) and value.startswith("="):
            return value
        return None
    finally:
        wb.close()
